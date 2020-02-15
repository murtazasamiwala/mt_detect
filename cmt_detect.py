"""Module to check if a translation matches with google translate."""
import os
from os.path import abspath
from googletrans import Translator
from difflib import SequenceMatcher
import re
import shutil
from zipfile import ZipFile
import win32com.client as win32
import xlrd
import pptx
import csv
import sys
from datetime import datetime
from google.oauth2 import service_account
from google.cloud import translate_v2 as translate
base_path = os.path.dirname(abspath('__file__'))
if 'dont_delete_ignore' not in os.listdir(base_path):
    os.mkdir('dont_delete_ignore')
    kmsg_1 = 'Key not found.\n'
    kmsg_2 = 'Key folder has been created.\nSave key file to this folder.'
    fkmsg = kmsg_1 + kmsg_2
    result = open('script_result.txt', 'w', encoding='utf8')
    result.write(fkmsg)
    result.close()
    sys.exit()
key_folder = base_path + '\\' + 'dont_delete_ignore'
key_path = key_folder + '\\' + os.listdir(key_folder)[0]
credentials = service_account.Credentials.from_service_account_file(
    key_path, scopes=["https://www.googleapis.com/auth/cloud-platform"])
ignored_fol = ['result_dir', 'dont_delete_ignore']
results_path = base_path + '\\' + 'result_dir'
csv_path = results_path + '\\' + 'results.csv'
excluded_files = ['_covering_letter.doc', '_Letter_from_the_Editor.docx']


def get_jc(path=base_path):
    """Get job code from zip or folders."""
    jc = None
    if jc is None:
        for i in os.listdir(path):
            parts = i.split('_')
            dir_condition = (os.path.isdir(i)) & (i not in ignored_fol)
            if (i.endswith('zip')) | (dir_condition):
                if parts[2] not in ['Original', 'Translate.zip', '604']:
                    jc = parts[0] + '_' + parts[1] + '_' + parts[2]
                else:
                    jc = parts[0] + '_' + parts[1]
    return jc


def extract_text(fname, path=base_path):
    """Extract text from given document."""
    if fname.split('.')[-1] in ['doc', 'docx', 'rtf']:
        word = win32.Dispatch('Word.Application')
        doc = word.Documents.Open(path+'\\'+fname)
        txt = doc.Content.Text
        doc.Close(False)
    elif fname.split('.')[-1] in ['xls', 'xlsx']:
        workbook = xlrd.open_workbook(fname)
        sheets_name = workbook.sheet_names()
        txt = '\n'
        for names in sheets_name:
            worksheet = workbook.sheet_by_name(names)
            num_rows = worksheet.nrows
            num_cells = worksheet.ncols
            for curr_row in range(num_rows):
                new_output = []
                for index_col in range(num_cells):
                    value = worksheet.cell_value(curr_row, index_col)
                    if value:
                        new_output.append(value)
                    if new_output:
                        txt += ' '.join(new_output) + '\n'
    elif fname.endswith('.pptx'):
        presentation = pptx.Presentation(fname)
        text_runs = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        txt = '\n\n'.join(text_runs)
    elif fname.endswith('.txt'):
        text_doc = open(fname, 'r', encoding='utf8')
        txt = text_doc.read()
        txt.close()
    elif fname.endswith('.csv'):
        csv_doc = open(fname, 'r', encoding='utf8')
        csv_reader = csv.reader(csv_doc, delimiter=',')
        txt = '\n'.join(['\t'.join(row) for row in csv_reader])
    return txt


def unzip_folders(path=base_path):
    """Unzip zip folders if any present."""
    for i in os.listdir(path):
        if i.endswith('.zip'):
            zf = ZipFile(i)
            zf.extractall()
    return


def folder_to_txt(path=base_path):
    """Extract text from document folders."""
    source = []
    translated = []
    source_names = []
    unzip_folders()
    for i in os.listdir(path):
        if (os.path.isdir(i)) & (i != 'results_dir'):
            for j in os.listdir(base_path + '\\' + i):
                if j == 'Job_files':
                    docs_path = base_path + '\\' + i + '\\' + 'Job_files'
                    for k in os.listdir(docs_path):
                        if k not in [jc + i for i in excluded_files]:
                            source.append(extract_text(k, docs_path))
                            source_names.append(k)
                        else:
                            pass
                elif j.endswith('Translate'):
                    docs_path = base_path + '\\' + i + '\\' + j
                    for k in os.listdir(docs_path):
                        if k not in [jc + i for i in excluded_files]:
                            translated.append(extract_text(k, docs_path))
    source_str = ''.join([text + '\n' for text in source])
    trans_str = ''.join([text + '\n'for text in translated])
    return source_str, trans_str, source_names


def detect_language(doc):
    """Detect language of given document."""
    lan = 'en'
    n = 0
    segs = len(doc) // 1000
    while lan == 'en':
        translator = Translator()
        if n < segs:
            lan = translator.detect(doc[n * 1000:(n + 1) * 1000]).lang
            n += 1
        else:
            lan = translator.detect(doc[n * 1000:]).lang
            break
    return lan


def doc_split(doc):
    """Split text into small chunks readable by google translate."""
    if language in ['ko', 'pt']:
        tokens = doc.split('.')
        tokens = [i + '.' for i in tokens]
    elif language in ['ja', 'zh-CN']:
        tokens = doc.split('。')
        tokens = [i + '。' for i in tokens]
    split = []
    len_counter = 0
    temp_list = []
    final = []
    for i in range(len(tokens)):
        if len_counter + len(tokens[i]) + len(temp_list) - 1 < 7000:
            len_counter = len_counter + len(tokens[i])
            temp_list.append(tokens[i])
        else:
            len_counter = len(tokens[i])
            split.append(temp_list)
            temp_list = []
            temp_list.append(tokens[i])
    split.append(temp_list)
    final = [''.join(i) for i in split]
    return final


def translate_text(split):
    """Check if text is already tranlsated. If not, translate it."""
    gt_list = []
    gt_out = None
    if 'result_dir' in os.listdir(base_path):
        if jc in os.listdir(results_path):
            gt_file = 'google_translated.txt'
            gt = open(job_path + gt_file, 'r', encoding='utf8')
            gt_out = gt.read()
            gt.close()
    if gt_out is None:
        for i in split:
            translate_client = translate.Client(credentials=credentials)
            result = translate_client.translate(i, target_language='en')
            gt_list.append(result['translatedText'])
        gt_out = ' '.join(gt_list)
    return gt_out


def save_files():
    """Save source, translated, and gt files in results folder."""
    if 'result_dir' not in os.listdir(base_path):
        os.mkdir('result_dir')
    if jc not in os.listdir(results_path):
        os.mkdir(results_path + '\\' + jc)
    for i in texts.keys():
        fname = open(job_path + '{}.txt'.format(i), 'w', encoding='utf8')
        fname.write(texts[i])
        fname.close()


def final_report():
    """Output results of test in txt and csv formats."""
    job_code = 'Job code: {}\n'.format(jc) + '*' * 20 + '\n'
    percent = round(ratio * 100, 2)
    sc = len(source)
    similarity = f'Translation is {percent}% similar to google translate\n'
    match_thou = round((high_matches / sc) * 1000)
    match_msg = f'{match_thou} long fragments per 1000 char match google\n'
    if ((high_matches / len(source)) > 3) | (percent > 35):
        decision = 'There seems high similarity to google. Please escalate'
    else:
        decision = 'Similarity is likely to be coincidental. Ignore'
    final_msg = job_code + similarity + match_msg + decision
    result = open('script_result.txt', 'w', encoding='utf8')
    result.write(final_msg)
    result.close()
    if 'results.csv' not in os.listdir(results_path):
        with open(csv_path, 'a', newline='') as result_csv:
            csv_writer = csv.writer(result_csv, delimiter=',')
            fields = ['job_code', 'date_time', 'source_chars',
                      'match_segments', 'percent_match', 'percent_segment',
                      'percent_length_high']
            csv_writer.writerow(fields)
            result_csv.close()
    dt = datetime.now().strftime("%d/%m/%Y %H:%M")
    pm = percent
    psm = round((high_matches / len(matches)) * 100, 2)
    test_doc_length = len(google_translated) + len(translated)
    plhm = round(((2 * len_high_matches)/test_doc_length) * 100, 2)
    result_list = [jc, dt, sc, high_matches, pm, psm, plhm]
    with open(csv_path, 'a', newline='') as result_csv:
        csv_writer = csv.writer(result_csv, delimiter=',')
        csv_writer.writerow(result_list)
        result_csv.close()


jc = get_jc()
source, translated, source_file_names = folder_to_txt()
similarity = SequenceMatcher(None, source, translated)
rep = 0
for i in similarity.get_matching_blocks():
    if i[2] > 30:
        if i[2] == len(source):
            pass
        elif source[i[0] - rep] in ['.', ' ']:
            rep_source = source[i[0] + 1 - rep:i[0] + i[2] - rep]
            source = source.replace(rep_source, '', 1)
            rep_trans = translated[i[1] + 1 - rep:i[1] + i[2] - rep]
            translated = translated.replace(rep_trans, '', 1)
            rep += i[2] - 1
        else:
            rep_source_2 = source[i[0] - rep:i[0] + i[2] - rep]
            source = source.replace(rep_source_2, '', 1)
            rep_trans_2 = translated[i[1] - rep:i[1] + i[2] - rep]
            translated = translated.replace(rep_trans_2, '', 1)
            rep += i[2]
job_path = results_path + '\\' + jc + '\\'
language = detect_language(source)
split_source = doc_split(source)
google_translated = translate_text(split_source)
zip_check = ['zip' in i for i in os.listdir(base_path)]
if any(zip_check):
    for i in os.listdir(base_path):
        if (os.path.isdir(i)) & (i not in ignored_fol):
            shutil.rmtree(i)
similarity2 = SequenceMatcher(None, google_translated, translated)
matches = similarity2.get_matching_blocks()
ratio = similarity2.ratio()
texts = {'google_translated': google_translated, 'translated': translated,
         'source': source}
save_files()
results = open('all_matches.txt', 'a', encoding='utf8')
high_matches = 0
len_high_matches = 0
match_threshold = {'pt': 100, 'ko': 80, 'ja': 80, 'zh-CN': 80}
for i in matches:
    if i[2] > match_threshold[language]:
        buffer = '\n' + '*' * 20 + '\n'
        results.write(google_translated[i[0]:i[0]+i[2]] + buffer)
        high_matches += 1
        len_high_matches += i[2]
results.close()
final_report()
