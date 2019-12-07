"""Module to check if a translation matches with google translate."""
import os
from os.path import abspath
from googletrans import Translator
from difflib import SequenceMatcher
import re
import shutil
from zipfile import ZipFile
import win32com.client as win32
import xlrd
import pptx
import csv
import sys
from datetime import datetime
base_path = os.path.dirname(abspath('__file__'))


def get_jc(path=base_path):
    """Get job code from zip or folders."""
    jc = None
    if jc is None:
        for i in os.listdir(path):
            parts = i.split('_')
            dir_condition = (os.path.isdir(i)) & (i != 'result_dir')
            if (i.endswith('zip')) | (dir_condition):
                if parts[2] not in ['Original', 'Translate.zip', '604']:
                    jc = parts[0] + '_' + parts[1] + '_' + parts[2]
                else:
                    jc = parts[0] + '_' + parts[1]
    return jc


def extract_text(fname, path=base_path):
    """Extract text from given document."""
    if fname.split('.')[-1] in ['doc', 'docx', 'rtf']:
        word = win32.Dispatch('Word.Application')
        doc = word.Documents.Open(path+'\\'+fname)
        txt = doc.Content.Text
        doc.Close(False)
        word.Quit()
    elif fname.split('.')[-1] in ['xls', 'xlsx']:
        workbook = xlrd.open_workbook(fname)
        sheets_name = workbook.sheet_names()
        txt = '\n'
        for names in sheets_name:
            worksheet = workbook.sheet_by_name(names)
            num_rows = worksheet.nrows
            num_cells = worksheet.ncols
            for curr_row in range(num_rows):
                new_output = []
                for index_col in range(num_cells):
                    value = worksheet.cell_value(curr_row, index_col)
                    if value:
                        new_output.append(value)
                    if new_output:
                        txt += ' '.join(new_output) + '\n'
    elif fname.endswith('.pptx'):
        presentation = pptx.Presentation(fname)
        text_runs = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        txt = '\n\n'.join(text_runs)
    elif fname.endswith('.txt'):
        text_doc = open(fname, 'r', encoding='utf8')
        txt = text_doc.read()
    elif fname.endswith('.csv'):
        csv_doc = open(fname, 'r', encoding='utf8')
        csv_reader = csv.reader(csv_doc, delimiter=',')
        txt = '\n'.join(['\t'.join(row) for row in csv_reader])
    return txt


def unzip_folders(path=base_path):
    """Unzip zip folders if any present."""
    for i in os.listdir(path):
        if i.endswith('.zip'):
            zf = ZipFile(i)
            zf.extractall()
    return


def folder_to_txt(path=base_path):
    """Extract text from document folders."""
    source = []
    translated = []
    source_names = []
    unzip_folders()
    for i in os.listdir(path):
        if (os.path.isdir(i)) & (i != 'results_dir'):
            for j in os.listdir(base_path + '\\' + i):
                if j == 'Job_files':
                    docs_path = base_path + '\\' + i + '\\' + 'Job_files'
                    for k in os.listdir(docs_path):
                        source.append(extract_text(k, docs_path))
                        source_names.append(k)
                elif j.endswith('Translate'):
                    docs_path = base_path + '\\' + i + '\\' + j
                    for k in os.listdir(docs_path):
                        translated.append(extract_text(k, docs_path))
    source_str = ''.join([text + '\n' for text in source])
    trans_str = ''.join([text + '\n'for text in translated])
    return source_str, trans_str, source_names


source, translated, source_file_names = folder_to_txt()
similarity = SequenceMatcher(None, source, translated)
rep = 0
for i in similarity.get_matching_blocks():
    if i[2] > 30:
        if i[2] == len(source):
            pass
        elif source[i[0] - rep] in ['.', ' ']:
            rep_source = source[i[0] + 1 - rep:i[0] + i[2] - rep]
            source = source.replace(rep_source, '', 1)
            rep_trans = translated[i[1] + 1 - rep:i[1] + i[2] - rep]
            translated = translated.replace(rep_trans, '', 1)
            rep += i[2] - 1
        else:
            rep_source_2 = source[i[0] - rep:i[0] + i[2] - rep]
            source = source.replace(rep_source_2, '', 1)
            rep_trans_2 = translated[i[1] - rep:i[1] + i[2] - rep]
            translated = translated.replace(rep_trans_2, '', 1)
            rep += i[2]
languages = {'en': 'english', 'pt': 'portuguese',
             'ko': 'korean', 'ja': 'japanese', 'zh-CN': 'chinese'}


def doc_split(doc):
    """Split text into small chunks readable by google translate."""
    translator = Translator()
    lan = translator.detect(doc[:1000]).lang
    if lan in ['ko', 'pt']:
        tokens = doc.split('.')
        tokens = [i + '.' for i in tokens]
    elif lan in ['ja', 'zh-CN']:
        tokens = doc.split('。')
        tokens = [i + '。' for i in tokens]
    split = []
    len_counter = 0
    temp_list = []
    final = []
    for i in range(len(tokens)):
        if len_counter + len(tokens[i]) + len(temp_list) - 1 < 1800:
            len_counter = len_counter + len(tokens[i])
            temp_list.append(tokens[i])
        else:
            len_counter = len(tokens[i])
            split.append(temp_list)
            temp_list = []
            temp_list.append(tokens[i])
    split.append(temp_list)
    final = [''.join(i) for i in split]
    return final, lan


google_output = []
split_source, language = doc_split(source)
for i in split_source:
    try:
        translator = Translator()
        google_output.append(translator.translate(i, dest='en').text)
    except ValueError:
        job_code = 'Job code: {}\n'.format(get_jc()) + '*' * 20 + '\n'
        msg_2 = 'Execution failed because of unrecognizabe characters.\n'
        msg_3 = 'Follow instructions for execution fail.'
        fin_msg = job_code + msg_2 + msg_3
        result = open('script_result.txt', 'w', encoding='utf-8')
        result.write(fin_msg)
        result.close()
        for i in os.listdir(base_path):
            if i.endswith('.zip'):
                os.remove(i)
        sys.exit()

zip_check = ['zip' in i for i in os.listdir(base_path)]
if any(zip_check):
    for i in os.listdir(base_path):
        if (os.path.isdir(i)) & (i != 'result_dir'):
            shutil.rmtree(i)
google_translated = ' '.join(google_output)
similarity2 = SequenceMatcher(None, google_translated, translated)
matches = similarity2.get_matching_blocks()
ratio = similarity2.ratio()
texts = {'google_translated': google_translated, 'translated': translated}
for i in texts.keys():
    fname = open('{}.txt'.format(i), 'w', encoding='utf8')
    fname.write(texts[i])
    fname.close()
results = open('all_matches.txt', 'a', encoding='utf8')
high_matches = 0
match_threshold = {'pt': 100, 'ko': 80, 'ja': 80, 'zh-CN': 80}
for i in matches:
    if i[2] > match_threshold[language]:
        results.write(google_translated[i[0]:i[0]+i[2]]+'\n')
        high_matches += 1
results.close()


def final_report():
    """Output results of test in txt and csv formats."""
    job_code = 'Job code: {}\n'.format(get_jc()) + '*' * 20 + '\n'
    percent = round(ratio * 100, 2)
    similarity = f'Translation is {percent}% similar to google translate\n'
    matches = f'{high_matches} fragments significantly match google\n'
    if (high_matches > 5) | (percent > 25):
        decision = 'There seems high similarity to google. Please escalate'
    else:
        decision = 'Similarity is likely to be coincidental. Ignore'
    final_msg = job_code + similarity + matches + decision
    result = open('script_result.txt', 'w', encoding='utf8')
    result.write(final_msg)
    result.close()
    if 'result_dir' not in os.listdir(base_path):
        os.mkdir('result_dir')
    results_path = base_path + '\\' + 'result_dir'
    csv_path = results_path + '\\' + 'results.csv'
    if 'results.csv' not in os.listdir(results_path):
        with open(csv_path, 'a', newline='') as result_csv:
            csv_writer = csv.writer(result_csv, delimiter=',')
            fields = ['job_code', 'date_time', 'source_chars',
                      'match_segments', 'percent_match']
            csv_writer.writerow(fields)
            result_csv.close()
    jc = get_jc()
    dt = datetime.now().strftime("%d/%m/%Y %H:%M")
    sc = len(source)
    ms = high_matches
    pm = percent
    result_list = [jc, dt, sc, ms, pm]
    with open(csv_path, 'a', newline='') as result_csv:
        csv_writer = csv.writer(result_csv, delimiter=',')
        csv_writer.writerow(result_list)
        result_csv.close()
final_report()
